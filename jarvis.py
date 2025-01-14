# pip install pyaudio

import pyttsx3 #pip install pyttsx3
import speech_recognition as sr #pip install speechRecognition
import datetime
import wikipedia #pip install wikipedia
import webbrowser
import os
import smtplib
from pptx import Presentation  # Change to this simpler import
import requests
from bs4 import BeautifulSoup  # Add these two imports
from colorama import init, Fore, Back, Style, Cursor
import time
import sys
import cv2
from pyzbar.pyzbar import decode
import pytesseract
pytesseract.pytesseract.tesseract_cmd = r'"C:\Users\ADMIN\Downloads\tesseract-ocr-w64-setup-5.5.0.20241111.exe"'  # You'll need to install Tesseract-OCR
import random
import re
from math import *  # For advanced mathematical operations
from geopy.geocoders import Nominatim
from geopy.distance import geodesic
from fer import FER
import matplotlib.pyplot as plt
from random import choice
import pygame  # pip install pygame
import tensorflow as tf
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
import spotipy  # pip install spotipy
from spotipy.oauth2 import SpotifyOAuth

# Add these lines to suppress TensorFlow messages
os.environ['TF_CPP_MIN_LOG_LEVEL'] = '2'
os.environ['TF_ENABLE_ONEDNN_OPTS'] = '0'

engine = pyttsx3.init('sapi5')
voices = engine.getProperty('voices')
# print(voices[1].id)
engine.setProperty('voice', voices[0].id)
engine.setProperty('rate', 180)
engine.setProperty('volume', 0.9)

# Spotify credentials setup
SPOTIPY_CLIENT_ID = '18b53a4379024135b898562d8db1d5f3'
SPOTIPY_CLIENT_SECRET = '27ec9410b0304380a16df36ea731247a'
SPOTIPY_REDIRECT_URI = 'http://localhost:8888/callback'

# Spotify Authorization
scope = 'user-read-playback-state user-modify-playback-state user-read-currently-playing'
sp = spotipy.Spotify(auth_manager=SpotifyOAuth(client_id=SPOTIPY_CLIENT_ID,
                                               client_secret=SPOTIPY_CLIENT_SECRET,
                                               redirect_uri=SPOTIPY_REDIRECT_URI,
                                               scope=scope))


def speak(audio):
    engine.say(audio)
    engine.runAndWait()


def wishMe():
    hour = int(datetime.datetime.now().hour)
    if hour>=0 and hour<12:
        speak("Good Morning!")

    elif hour>=12 and hour<18:
        speak("Good Afternoon!")   

    else:
        speak("Good Evening!")  

    speak("I am jarvis sir how can i help you. ")       

def takeCommand():
    r = sr.Recognizer()
    with sr.Microphone() as source:
        print("Listening...")
        # Optimized settings for English recognition
        r.pause_threshold = 0.5
        r.energy_threshold = 4000  # Increased for better clarity
        r.dynamic_energy_threshold = True
        r.adjust_for_ambient_noise(source, duration=0.2)
        
        try:
            print("Speak now...")
            audio = r.listen(source, timeout=5, phrase_time_limit=5)
            print("Processing...")
        except sr.WaitTimeoutError:
            print("Listening timed out")
            return "None"

    try:
        print("Recognizing...")    
        # Force English US recognition only
        query = r.recognize_google(audio, language='en-US')
        print(f"User said: {query}\n")
        return query.lower()

    except sr.UnknownValueError:
        print("Could not understand audio. Please speak in English.")
        return "None"
    except sr.RequestError:
        print("Service error. Check your internet connection.")
        return "None"
    except Exception as e:
        print("Error occurred. Please try again.")
        return "None"
# Spotify credentials setup
SPOTIPY_CLIENT_ID = 'your_spotify_client_id'
SPOTIPY_CLIENT_SECRET = 'your_spotify_client_secret'
SPOTIPY_REDIRECT_URI = 'http://localhost:8888/callback'

# Spotify Authorization
scope = 'user-read-playback-state user-modify-playback-state user-read-currently-playing'
sp = spotipy.Spotify(auth_manager=SpotifyOAuth(client_id=SPOTIPY_CLIENT_ID,
                                               client_secret=SPOTIPY_CLIENT_SECRET,
                                               redirect_uri=SPOTIPY_REDIRECT_URI,
                                               scope=scope))
def sendEmail(to_email=None, subject=None, content=None):
    try:
        # Your email credentials
        EMAIL_ADDRESS = "your_email@gmail.com"  # Replace with your email
        EMAIL_PASSWORD = "your_app_password"    # Use App Password from Google Account
        
        # Validate inputs
        if not all([to_email, subject, content]):
            speak("Missing email information")
            return False
            
        # Validate email format
        if '@' not in to_email or '.' not in to_email:
            speak("Invalid email format")
            return False
        
        # Create message
        msg = MIMEMultipart()
        msg['From'] = EMAIL_ADDRESS
        msg['To'] = to_email
        msg['Subject'] = subject
        msg.attach(MIMEText(content, 'plain'))
        
        # Create server connection
        server = smtplib.SMTP('smtp.gmail.com', 587)
        server.starttls()
        
        try:
            server.login(EMAIL_ADDRESS, EMAIL_PASSWORD)
            text = msg.as_string()
            server.sendmail(EMAIL_ADDRESS, to_email, text)
            return True
            
        except smtplib.SMTPAuthenticationError:
            speak("Email authentication failed. Please check your credentials.")
            return False
        except smtplib.SMTPException as e:
            speak(f"SMTP error occurred: {str(e)}")
            return False
        finally:
            server.quit()
            
    except Exception as e:
        print(f"Email error: {str(e)}")
        return False

def clear_screen():
    os.system('cls' if os.name == 'nt' else 'clear')

def print_centered(text):
    terminal_width = os.get_terminal_size().columns
    print(text.center(terminal_width))

def startup_animation():
    init()  # Initialize colorama
    clear_screen()
    
    # Matrix-style rain animation
    def matrix_rain():
        chars = "10"
        for _ in range(3):  # Number of rain cycles
            for i in range(10):  # Rain drops
                print(Fore.GREEN + choice(chars), end='')
                sys.stdout.flush()
                time.sleep(0.05)
            print()
    
    # JARVIS Logo with color transition
    logo = """
    ╔══════════════════════════════════════════════╗
     ██╗ █████╗ ██████╗ ██╗   ██╗██╗███████╗
     ██║██╔══██╗██╔══██╗██║   ██║██║██╔════╝
     ██║███████║██████╔╝██║   ██║██║███████╗
██   ██║██╔══██║██╔══██╗╚██╗ ██╔╝██║╚════██║
╚█████╔╝██║  ██║██║  ██║ ╚████╔╝ ██║███████║
 ╚════╝ ╚═╝  ╚═╝╚═╝  ╚═╝  ╚═══╝  ╚═╝╚══════╝
    ╚══════════════════════════════════════════════╝
    """
    
    # Animated loading sequence
    loading_chars = ["⠋", "⠙", "⠹", "⠸", "⠼", "⠴", "⠦", "⠧", "⠇", "⠏"]
    
    # Matrix rain effect
    matrix_rain()
    time.sleep(0.5)
    clear_screen()
    
    # Animated logo reveal
    for line in logo.split('\n'):
        print_centered(Fore.CYAN + line)
        time.sleep(0.1)
    
    # System initialization text
    init_texts = [
        "INITIALIZING SYSTEMS",
        "LOADING AI CORE",
        "CALIBRATING VOICE SYSTEMS",
        "ESTABLISHING NEURAL NETWORKS",
        "OPTIMIZING RESPONSE ALGORITHMS"
    ]
    
    print("\n")
    for text in init_texts:
        # Spinning animation for each initialization step
        for char in loading_chars:
            sys.stdout.write(f"\r{Fore.YELLOW}{char} {text}...")
            sys.stdout.flush()
            time.sleep(0.1)
        print(Fore.GREEN + " DONE")
        time.sleep(0.2)
    
    # Progress bar animation
    print("\n" + Fore.CYAN + "Final System Check:" + Style.RESET_ALL)
    for i in range(21):
        sys.stdout.write('\r')
        sys.stdout.write(Fore.GREEN + "[%-20s] %d%%" % ('='*i, 5*i))
        sys.stdout.flush()
        time.sleep(0.1)
    
    print("\n\n" + Fore.LIGHTBLUE_EX + "✓ All systems operational!" + Style.RESET_ALL)
    time.sleep(1)
    
    # Pulse effect for ready message
    for _ in range(3):
        clear_screen()
        print(logo)
        print_centered(Fore.GREEN + "JARVIS READY" + Style.RESET_ALL)
        time.sleep(0.3)
        clear_screen()
        print(logo)
        print_centered(Fore.LIGHTGREEN_EX + "JARVIS READY" + Style.RESET_ALL)
        time.sleep(0.3)

def scan_camera():
    try:
        cap = cv2.VideoCapture(0)  # Initialize camera
        speak("Starting advanced scan mode")
        print("Controls:")
        print("- Press 's' to scan")
        print("- Press 'q' to quit")
        speak("Press s to scan or q to quit")
        
        while True:
            ret, frame = cap.read()
            if not ret:
                speak("Cannot access camera")
                break
            
            # Add scanning frame overlay
            height, width = frame.shape[:2]
            cv2.rectangle(frame, (width//4, height//4), 
                         (3*width//4, 3*height//4), (0, 255, 0), 2)
            
            # Add instructions text
            cv2.putText(frame, "Press 's' to scan, 'q' to quit", 
                       (10, 30), cv2.FONT_HERSHEY_SIMPLEX, 
                       0.7, (0, 255, 0), 2)
                
            # Display the frame
            cv2.imshow('Jarvis Scanner', frame)
            
            key = cv2.waitKey(1)
            if key & 0xFF == ord('q'):  # Quit
                speak("Exiting scanner mode")
                break
            elif key & 0xFF == ord('s'):  # Scan
                speak("Scanning image")
                
                # Crop to scanning area
                scan_area = frame[height//4:3*height//4, width//4:3*width//4]
                
                # Try to decode QR codes
                qr_results = decode(scan_area)
                if qr_results:
                    for qr in qr_results:
                        qr_text = qr.data.decode('utf-8')
                        speak(f"QR Code detected")
                        print(f"QR Content: {qr_text}")
                        speak(f"Content is: {qr_text}")
                
                # Try to read text
                # Convert to grayscale for better OCR
                gray = cv2.cvtColor(scan_area, cv2.COLOR_BGR2GRAY)
                # Apply thresholding to get better text recognition
                _, threshold = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
                
                # Configure tesseract parameters
                custom_config = r'--oem 3 --psm 6'
                text = pytesseract.image_to_string(threshold, config=custom_config)
                
                if text.strip():
                    print("\nDetected Text:")
                    print("-" * 50)
                    print(text.strip())
                    print("-" * 50)
                    speak("I found some text in the image")
                    
                    # Save the scanned text to a file
                    timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
                    filename = f"scan_{timestamp}"
                    
                    # Save image
                    cv2.imwrite(f"{filename}.jpg", scan_area)
                    # Save text
                    with open(f"{filename}.txt", "w") as f:
                        f.write(text)
                    
                    speak("I've saved both the image and detected text to files")
                    print(f"Files saved as {filename}.jpg and {filename}.txt")
                else:
                    speak("No text was detected in the scan area")
        
        cap.release()
        cv2.destroyAllWindows()
        
    except Exception as e:
        speak("An error occurred with the scanner")
        print(f"Error: {e}")
        
def take_picture():
    try:
        cap = cv2.VideoCapture(0)  # Initialize camera
        speak("Camera activated. Press 'spacebar' to take picture, 'q' to quit.")
        
        while True:
            ret, frame = cap.read()
            if not ret:
                speak("Cannot access camera")
                break
                
            # Display the frame
            cv2.imshow('Jarvis Camera', frame)
            
            key = cv2.waitKey(1)
            if key & 0xFF == ord('q'):  # Quit
                break
            elif key == 32:  # Spacebar
                # Generate timestamp for unique filename
                timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
                filename = f"picture_{timestamp}.jpg"
                
                # Save the picture
                cv2.imwrite(filename, frame)
                speak(f"Picture taken and saved as {filename}")
                print(f"Picture saved as: {filename}")
        
        cap.release()
        cv2.destroyAllWindows()
        
    except Exception as e:
        speak("An error occurred with the camera")
        print(e)
        
def get_response(query):
    # Basic responses to common queries
    responses = {
        'how are you': ['I am doing well, thank you for asking!', 'I am great! How are you?', 'All systems operational!'],
        'what is your name': ['I am Jarvis, your AI assistant.', 'My name is Jarvis, nice to meet you!'],
        'who made you': ['I was created by my developer as an AI assistant.', 'I am an AI assistant created to help you.'],
        'what can you do': ['I can help you with various tasks like opening applications, searching the web, taking pictures, and having conversations!'],
        'tell me a joke': ['Why don\'t programmers like nature? It has too many bugs!', 
                          'What do you call a bear with no teeth? A gummy bear!',
                          'Why did the AI go to therapy? It had too many processing issues!'],
        'hello': ['Hello! How can I help you today?', 'Hi there! What can I do for you?'],
        'kutte': [ 'tu kutta tera pura khandan kutta'],
        'good morning': ['Good morning! Hope you have a great day!', 'Good morning! How can I assist you today?'],
        'good afternoon': ['Good afternoon! How may I help you?', 'Good afternoon! What can I do for you?'],
        'good evening': ['Good evening! How can I assist you?', 'Good evening! What would you like me to do?'],
        'thank you': ['You\'re welcome!', 'Glad I could help!', 'My pleasure!'],
        'how old are you': ['I am an AI, so I don\'t age like humans do!', 'Age is just a number for AI assistants like me!'],
        'are you real': ['I am an AI assistant, designed to help make your life easier!', 'I\'m as real as software can be!']
    }
    
    # Check for matching queries
    for key in responses:
        if key in query.lower():
            return random.choice(responses[key])
    
    return None

def calculate_expression(expression):
    try:
        # Remove any text that's not part of the mathematical expression
        expression = re.sub(r'[^0-9+\-*/().%\s]', '', expression)
        
        # Evaluate the expression
        result = eval(expression)
        return str(result)
    except:
        return "Sorry, I couldn't calculate that."

def get_location_info(place):
    try:
        geolocator = Nominatim(user_agent="jarvis_assistant")
        location = geolocator.geocode(place, timeout=10)
        
        if location:
            address = location.address
            latitude = location.latitude
            longitude = location.longitude
            return address, (latitude, longitude)
        return None, None
    except Exception as e:
        print(f"Error getting location: {e}")
        return None, None

def calculate_distance(place1, place2):
    try:
        loc1_info, loc1_coords = get_location_info(place1)
        loc2_info, loc2_coords = get_location_info(place2)
        
        if loc1_coords and loc2_coords:
            distance = geodesic(loc1_coords, loc2_coords).kilometers
            return distance
        return None
    except Exception as e:
        print(f"Error calculating distance: {e}")
        return None

def analyze_expression():
    try:
        # Initialize camera and detector
        cap = cv2.VideoCapture(0)
        detector = FER(mtcnn=True)
        speak("Starting facial expression analysis. Press 'q' to quit, 's' to save analysis.")
        
        # Variables for emotion tracking
        emotions_history = []
        
        while True:
            ret, frame = cap.read()
            if not ret:
                speak("Cannot access camera")
                break
            
            # Detect emotions in the frame
            result = detector.detect_emotions(frame)
            
            # Process and display results
            if result:
                # Get the first face detected
                emotions = result[0]['emotions']
                dominant_emotion = max(emotions.items(), key=lambda x: x[1])[0]
                confidence = emotions[dominant_emotion]
                
                # Store emotion data
                emotions_history.append(emotions)
                
                # Draw rectangle around face
                box = result[0]['box']
                cv2.rectangle(frame,
                            (box[0], box[1]),
                            (box[0] + box[2], box[1] + box[3]),
                            (0, 255, 0), 2)
                
                # Display emotion text with better formatting
                y_offset = box[1] - 10
                for emotion, score in emotions.items():
                    text = f"{emotion}: {score:.2f}"
                    cv2.putText(frame, text, 
                              (box[0], y_offset),
                              cv2.FONT_HERSHEY_SIMPLEX, 0.6,
                              (0, 255, 0), 2)
                    y_offset -= 20
                
                # Speak dominant emotion if it changes significantly
                if len(emotions_history) > 1:
                    prev_dominant = max(emotions_history[-2].items(), key=lambda x: x[1])[0]
                    if dominant_emotion != prev_dominant:
                        speak(f"Detected {dominant_emotion}")
            
            # Add instructions overlay
            cv2.putText(frame, "Press 'q' to quit, 's' to save analysis",
                       (10, 30), cv2.FONT_HERSHEY_SIMPLEX, 0.7,
                       (0, 255, 0), 2)
            
            # Display the frame
            cv2.imshow('Expression Analysis', frame)
            
            key = cv2.waitKey(1) & 0xFF
            if key == ord('q'):
                break
            elif key == ord('s'):
                # Save analysis results
                timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
                
                # Save the current frame
                cv2.imwrite(f"expression_{timestamp}.jpg", frame)
                
                # Create emotion analysis graph
                if emotions_history:
                    plt.figure(figsize=(10, 6))
                    avg_emotions = {emotion: sum(hist[emotion] for hist in emotions_history) / len(emotions_history)
                                  for emotion in emotions_history[0].keys()}
                    
                    plt.bar(avg_emotions.keys(), avg_emotions.values())
                    plt.title('Average Emotion Analysis')
                    plt.ylabel('Confidence Score')
                    plt.xticks(rotation=45)
                    plt.tight_layout()
                    plt.savefig(f"emotion_analysis_{timestamp}.png")
                    plt.close()
                    
                    speak("Analysis saved successfully")
        
        cap.release()
        cv2.destroyAllWindows()
        
        # Final analysis
        if emotions_history:
            speak("Analysis complete. Here's the summary.")
            avg_emotions = {emotion: sum(hist[emotion] for hist in emotions_history) / len(emotions_history)
                          for emotion in emotions_history[0].keys()}
            dominant_overall = max(avg_emotions.items(), key=lambda x: x[1])[0]
            speak(f"The dominant emotion was {dominant_overall}")
        
    except Exception as e:
        speak("An error occurred with expression analysis")
        print(f"Error: {e}")
        if 'cap' in locals():
            cap.release()
        cv2.destroyAllWindows()

def play_audio(file_path):
    try:
        pygame.mixer.init()
        pygame.mixer.music.load(file_path)
        pygame.mixer.music.play()
        while pygame.mixer.music.get_busy():
            pygame.time.Clock().tick(10)
    except Exception as e:
        print(f"Error playing audio: {e}")
        speak("Sorry, I couldn't play that audio file")

def write_code(language, description):
    try:
        # Basic code templates
        templates = {
            'python': {
                'hello': '''print("Hello, World!")''',
                
                'function': '''def my_function():
    # Function code here
    pass''',
                
                'class': '''class MyClass:
    def __init__(self):
        pass
    
    def my_method(self):
        pass''',
                
                'loop': '''for i in range(10):
    print(i)''',
                
                'if': '''if condition:
    # Code here
    pass
else:
    # Code here
    pass''',
                
                'file': '''with open('filename.txt', 'r') as file:
    content = file.read()
    print(content)'''
            },
            
            'html': {
                'basic': '''<!DOCTYPE html>
<html>
<head>
    <title>My Page</title>
</head>
<body>
    <h1>Hello World</h1>
</body>
</html>'''
            }
        }
        
        timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        
        if 'python' in language.lower():
            filename = f"generated_code_{timestamp}.py"
            extension = '.py'
        elif 'html' in language.lower():
            filename = f"generated_code_{timestamp}.html"
            extension = '.html'
        else:
            filename = f"generated_code_{timestamp}.txt"
            extension = '.txt'
            
        # Generate code based on description
        code = ""
        if 'hello' in description.lower():
            code = templates[language]['hello']
        elif 'function' in description.lower():
            code = templates[language]['function']
        elif 'class' in description.lower():
            code = templates[language]['class']
        elif 'loop' in description.lower():
            code = templates[language]['loop']
        elif 'file' in description.lower():
            code = templates[language]['file']
        else:
            code = templates[language]['basic']
            
        # Save the code to a file
        with open(filename, 'w') as file:
            file.write(code)
            
        speak(f"I've generated the code and saved it as {filename}")
        print(f"Generated code saved in {filename}:")
        print("-" * 50)
        print(code)
        print("-" * 50)
        
        return True
        
    except Exception as e:
        print(f"Error generating code: {e}")
        speak("Sorry, I couldn't generate the code")
        return False

if __name__ == "__main__":
    startup_animation()
    wishMe()
    while True:
        query = takeCommand()
        if query == "None":
            continue

        # Logic for executing tasks based on query
        if 'wikipedia' in query:
            speak('Searching Wikipedia...')
            query = query.replace("wikipedia", "")
            results = wikipedia.summary(query, sentences=2)
            speak("According to Wikipedia")
            print(results)
            speak(results)

        elif 'open youtube' in query:
            webbrowser.open("youtube.com")

        elif 'open google' in query:
            webbrowser.open("google.com")

        elif 'open physics' in query:
            webbrowser.open("https://www.youtube.com/@PhysicsWallah")   


        elif 'play music' in query:
            music_dir = 'C:\\Users\\ADMIN\\Music\\New folder'
            songs = os.listdir(music_dir)
            print(songs)    
            os.startfile(os.path.join(music_dir, songs[0]))

        elif 'next song' in query:
            music_dir = 'C:\\Users\\ADMIN\\Music\\New folder'
            songs = os.listdir(music_dir)
            current_song = songs.index(current_playing) if 'current_playing' in locals() else -1
            next_song = (current_song + 1) % len(songs)
            current_playing = songs[next_song]
            os.startfile(os.path.join(music_dir, current_playing))
            speak(f"Playing next song: {current_playing}")

        elif 'previous song' in query:
            music_dir = 'C:\\Users\\ADMIN\\Music\\New folder'
            songs = os.listdir(music_dir)
            current_song = songs.index(current_playing) if 'current_playing' in locals() else 0
            prev_song = (current_song - 1) % len(songs)
            current_playing = songs[prev_song]
            os.startfile(os.path.join(music_dir, current_playing))
            speak(f"Playing previous song: {current_playing}")

        elif 'time' in query:
            strTime = datetime.datetime.now().strftime("%H:%M:%S")    
            speak(f"Sir, the time is {strTime}")

        elif 'open code' in query:
            codePath = "C:\\Users\\ADMIN\\AppData\\Roaming\\Microsoft\\Windows\\Start Menu\\Programs\\Visual Studio Code\\Visual Studio Code.lnk"
            os.startfile(codePath)

        elif 'make presentation' in query:
            try:
                speak("Creating a creative PowerPoint presentation")
                prs = Presentation()
                
                # Title slide
                title_slide = prs.slides.add_slide(prs.slide_layouts[0])
                title = title_slide.shapes.title
                subtitle = title_slide.placeholders[1]
                title.text = "Creative Presentation"
                subtitle.text = "Created by Jarvis AI Assistant"
                
                # Content slide 1 - Bullet points
                bullet_slide = prs.slides.add_slide(prs.slide_layouts[1])
                shapes = bullet_slide.shapes
                title_shape = shapes.title
                body_shape = shapes.placeholders[1]
                title_shape.text = "Key Points"
                
                tf = body_shape.text_frame
                tf.text = "Main Ideas:"
                tf.add_paragraph().text = "• Innovation and Creativity"
                tf.add_paragraph().text = "• Strategic Thinking"
                tf.add_paragraph().text = "• Future Vision"
                
                # Content slide 2 - Two content
                two_content_slide = prs.slides.add_slide(prs.slide_layouts[3])
                shapes = two_content_slide.shapes
                title_shape = shapes.title
                title_shape.text = "Comparison"
                
                # Left and right content
                body_shape = shapes.placeholders[1]
                tf = body_shape.text_frame
                tf.text = "Current State:"
                tf.add_paragraph().text = "• Present scenario"
                tf.add_paragraph().text = "• Existing solutions"
                
                other_body_shape = shapes.placeholders[2]
                tf = other_body_shape.text_frame
                tf.text = "Future Goals:"
                tf.add_paragraph().text = "• Upcoming developments"
                tf.add_paragraph().text = "• Expected outcomes"
                
                # Thank you slide
                thank_slide = prs.slides.add_slide(prs.slide_layouts[5])
                title_shape = thank_slide.shapes.title
                title_shape.text = "Thank You!"
                
                # Save the presentation
                prs.save('Creative_Presentation.pptx')
                speak("I've created a creative presentation with multiple slides including title, content, and thank you slides!")
                
            except Exception as e:
                print(e)
                speak("Sorry, I encountered an error while creating the presentation")

        elif 'turn off' in query or 'goodbye' in query or 'exit' in query:
            speak("Goodbye sir, have a nice day!")
            exit()

        elif 'google' in query:
            try:
                search_query = query.replace('google', '').strip()
                speak(f"Searching Google for {search_query}")
                search_url = f"https://www.google.com/search?q={search_query.replace(' ', '+')}"
                webbrowser.open(search_url)
                speak("I have opened the search results")
                
            except Exception as e:
                speak("Sorry, I couldn't perform the search")
                print(e)

        elif 'scan' in query or 'camera' in query:
            speak("Activating scanner mode")
            scan_camera()

        elif 'take picture' in query or 'take photo' in query or 'click photo' in query:
            speak("Opening camera for taking picture")
            take_picture()

        elif any(word in query for word in ['how are you', 'what is your name', 'who made you', 
                'what can you do', 'kutte', 'tell me a joke', 'hello', 'good morning', 'good afternoon', 
                'good evening', 'thank you', 'how old are you', 'are you real']):
            response = get_response(query)
            if response:
                speak(response)
                print(response)

        elif any(word in query for word in ['calculate', 'solve', 'what is', 'equals']):
            try:
                # Extract the mathematical expression
                expression = query.replace('calculate', '')
                expression = expression.replace('solve', '')
                expression = expression.replace('what is', '')
                expression = expression.replace('equals', '')
                expression = expression.strip()
                
                result = calculate_expression(expression)
                speak(f"The result is {result}")
                print(f"Result: {result}")
                
            except Exception as e:
                speak("Sorry, I couldn't solve that mathematical problem")
                print(e)

        elif 'add' in query or 'plus' in query:
            try:
                numbers = [int(num) for num in re.findall(r'\d+', query)]
                if len(numbers) >= 2:
                    result = sum(numbers)
                    speak(f"The sum is {result}")
                    print(f"Sum: {result}")
            except:
                speak("Sorry, I couldn't perform the addition")

        elif 'multiply' in query or 'times' in query or 'multiplication' in query:
            try:
                # Remove words that might interfere with number extraction
                query = query.replace('multiply', '').replace('times', '').replace('multiplication', '')
                query = query.replace('by', '').replace('and', '').replace('with', '')
                
                # Extract numbers including decimals
                numbers = [float(num) for num in re.findall(r'\d*\.?\d+', query)]
                
                if len(numbers) >= 2:
                    result = 1
                    for num in numbers:
                        result *= num
                    
                    # Format result to avoid unnecessary decimals
                    if result.is_integer():
                        result = int(result)
                    else:
                        result = round(result, 2)
                        
                    speak(f"The product is {result}")
                    print(f"Product: {result}")
                else:
                    speak("Please provide at least two numbers for multiplication")
            except Exception as e:
                speak("Sorry, I couldn't perform the multiplication")
                print(f"Error: {e}")

        elif 'find location' in query or 'where is' in query or 'locate' in query:
            try:
                place = query.replace('find location', '').replace('where is', '').replace('locate', '').strip()
                speak(f"Searching for {place}")
                
                address, coords = get_location_info(place)
                if address and coords:
                    speak(f"I found {place}")
                    print(f"Address: {address}")
                    print(f"Coordinates: {coords}")
                    
                    # Open in Google Maps
                    maps_url = f"https://www.google.com/maps/search/?api=1&query={coords[0]},{coords[1]}"
                    webbrowser.open(maps_url)
                    speak("I've opened the location in Google Maps")
                else:
                    speak(f"Sorry, I couldn't find {place}")
                    
            except Exception as e:
                speak("Sorry, I couldn't process that location request")
                print(e)

        elif 'distance between' in query or 'how far' in query:
            try:
                # Extract two locations from query
                query = query.replace('distance between', '').replace('how far is', '').replace('from', '').replace('to', '')
                places = query.split('and')
                if len(places) == 2:
                    place1 = places[0].strip()
                    place2 = places[1].strip()
                    
                    distance = calculate_distance(place1, place2)
                    if distance:
                        speak(f"The distance between {place1} and {place2} is approximately {round(distance, 2)} kilometers")
                        print(f"Distance: {round(distance, 2)} km")
                    else:
                        speak("Sorry, I couldn't calculate the distance between these locations")
                else:
                    speak("Please specify two locations to calculate distance")
                    
            except Exception as e:
                speak("Sorry, I couldn't calculate the distance")
                print(e)

        elif 'analyze expression' in query or 'check expression' in query or 'how do i look' in query:
            speak("Starting facial expression analysis")
            analyze_expression()

        elif 'send email' in query or 'send mail' in query:
            try:
                # Get email address
                speak("Who should I send the email to? Please spell out the email address")
                print("Format example: example at gmail dot com")
                to_email = takeCommand()
                while to_email == "None":
                    speak("I didn't catch that. Please spell out the email address again")
                    to_email = takeCommand()
                
                # Convert spoken email to proper format
                to_email = to_email.lower().replace(" at ", "@").replace(" dot ", ".").replace(" ", "")
                speak(f"Sending to {to_email}. Is this correct? Say yes to confirm.")
                confirm = takeCommand()
                if 'yes' not in confirm.lower():
                    speak("Email cancelled")
                    continue
                
                # Get subject
                speak("What should be the subject of the email?")
                subject = takeCommand()
                while subject == "None":
                    speak("I didn't catch that. Please say the subject again")
                    subject = takeCommand()
                print(f"Subject: {subject}")
                speak(f"Subject will be: {subject}")
                
                # Get content
                speak("What should I say in the email?")
                content = takeCommand()
                while content == "None":
                    speak("I didn't catch that. Please say the content again")
                    content = takeCommand()
                print(f"Content: {content}")
                speak(f"Content will be: {content}")
                
                # Final confirmation
                speak("Should I send this email? Say yes to confirm.")
                final_confirm = takeCommand()
                
                if 'yes' in final_confirm.lower():
                    if sendEmail(to_email, subject, content):
                        speak("Email has been sent successfully!")
                    else:
                        speak("Sorry, there was an error sending the email")
                else:
                    speak("Email cancelled")
                    
            except Exception as e:
                print(f"Email error: {e}")
                speak("Sorry, I couldn't process the email request")

        elif 'write code' in query or 'generate code' in query or 'create program' in query:
            try:
                speak("What programming language should I use?")
                language = takeCommand().lower()
                while language == "None":
                    speak("I didn't catch that. Please specify the programming language")
                    language = takeCommand().lower()
                
                speak("What kind of code would you like me to write?")
                description = takeCommand().lower()
                while description == "None":
                    speak("I didn't catch that. Please describe the code you want")
                    description = takeCommand().lower()
                
                if write_code(language, description):
                    speak("Code has been generated successfully!")
                else:
                    speak("There was an error generating the code")
                    
            except Exception as e:
                print(e)
                speak("Sorry, I couldn't process the code generation request")

        else:
            print("No query matched")

    print(sys.executable)  # Add this temporarily to check your Python path
    print(sys.path)  # This will show where Python looks for modules
